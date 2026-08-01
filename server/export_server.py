"""
AuraAI — Native PowerPoint Compiler v1.1
Converts AI-generated HTML slides into a real, editable .pptx file
using python-pptx + BeautifulSoup.

Production:  uvicorn export_server:app --host 0.0.0.0 --port $PORT
Development: uvicorn export_server:app --reload --port 4567
"""

import io
import os
import re
import time
import traceback
from collections import defaultdict

from bs4 import BeautifulSoup
from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.gzip import GZipMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pydantic import BaseModel

# ─── Safety Wrappers ──────────────────────────────────────────────────────────

FORBIDDEN_TAGS = re.compile(r'<(svg|canvas|iframe|script|object|embed|video|audio)[\s\S]*?</\1>', re.IGNORECASE)

def sanitize_html(html: str) -> str:
    """Strip tags that crash the Python XML parser."""
    return FORBIDDEN_TAGS.sub('', html).strip()

def safe_inch(val, min_val=0.1):
    """Guarantees dimensions are never zero or negative."""
    try:
        return Inches(max(float(val), min_val))
    except:
        return Inches(min_val)

def safe_hex(hex_str, default="FFFFFF"):
    """Prevents invalid color strings from crashing the XML."""
    try:
        if not hex_str: return parse_color(default)
        h = hex_str.strip().lstrip('#')
        if len(h) == 3: h = ''.join([c*2 for c in h])
        if len(h) != 6: h = default
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    except:
        return RGBColor(255,255,255)

def safe_font_size(pt, min_val=8, max_val=72):
    """Clamps font sizes to reasonable bounds."""
    try:
        return Pt(max(min(float(pt), max_val), min_val))
    except:
        return Pt(12)

# ─── Config ───────────────────────────────────────────────────────────────────

MAX_HTML_SIZE = int(os.getenv("PPTX_MAX_HTML_SIZE", 5_000_000))   # ~5 MB
RATE_LIMIT_PER_MIN = int(os.getenv("PPTX_RATE_LIMIT", 30))
CORS_ORIGINS = os.getenv("PPTX_CORS_ORIGINS", "*").split(",")
CACHE_SECONDS = int(os.getenv("PPTX_CACHE_SECONDS", 0))

app = FastAPI(title="AuraAI PPTX Compiler")

app.add_middleware(GZipMiddleware, minimum_size=500)
app.add_middleware(
    CORSMiddleware,
    allow_origins=CORS_ORIGINS,
    allow_methods=["POST", "GET"],
    allow_headers=["*"],
)

# ─── In-memory rate limiter ──────────────────────────────────────────────────

_rate_store: dict[str, list[float]] = defaultdict(list)


def _check_rate(ip: str):
    now = time.time()
    window = _rate_store[ip]
    window[:] = [t for t in window if now - t < 60]
    if len(window) >= RATE_LIMIT_PER_MIN:
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Try again in a minute.")
    window.append(now)

# ─── Theme fallback colors ───────────────────────────────────────────────────

THEME_DEFAULTS = {
    "glass-dark": {
        "bg": "0F0C29", "title": "FAFAFA", "body": "C7C9D9",
        "card_bg": "1E1B4B", "card_border": "3F3D6E", "accent": "06B6D4",
        "table_header_bg": "26224D", "table_row_bg": "17153A",
    },
    "minimal-light": {
        "bg": "FAFAFA", "title": "111827", "body": "4B5563",
        "card_bg": "FFFFFF", "card_border": "E5E7EB", "accent": "2563EB",
        "table_header_bg": "F3F4F6", "table_row_bg": "FFFFFF",
    },
    "neo-brutalism": {
        "bg": "FFEB3B", "title": "000000", "body": "1A1A1A",
        "card_bg": "FFFFFF", "card_border": "000000", "accent": "FF3D57",
        "table_header_bg": "FF3D57", "table_row_bg": "F5F5F5",
    },
}
DEFAULT_THEME = "glass-dark"


class ExportRequest(BaseModel):
    html: str
    theme: str = DEFAULT_THEME
    filename: str = "presentation.pptx"

# ─── Helpers ──────────────────────────────────────────────────────────────────

def extract_styles(el):
    styles = {}
    style_str = el.get("style", "") if el else ""
    if not style_str:
        return styles
    for pair in style_str.split(";"):
        if ":" in pair:
            k, v = pair.split(":", 1)
            styles[k.strip().lower()] = v.strip()
    return styles


def get_theme_colors(slide_el, fallback_theme):
    theme_id = slide_el.get("data-theme") or fallback_theme
    return THEME_DEFAULTS.get(theme_id, THEME_DEFAULTS[DEFAULT_THEME]), theme_id


def pct_from_style(style_str, prop):
    m = re.search(rf"{prop}\s*:\s*([\d.]+)%", style_str or "")
    return float(m.group(1)) if m else None


# ─── Compiler ─────────────────────────────────────────────────────────────────

def compile_html_to_pptx(html_content: str, default_theme: str, output: io.BytesIO):
    soup = BeautifulSoup(html_content, "html.parser")
    slides = soup.find_all("section")

    if not slides:
        raise ValueError("No <section> slides found in HTML")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for i, slide_html in enumerate(slides):
        try:
            colors, theme_id = get_theme_colors(slide_html, default_theme)
            slide = prs.slides.add_slide(blank_layout)
            s_styles = extract_styles(slide_html)

            # ── Background ──
            bg_hex = s_styles.get("background") or s_styles.get("background-color") or colors["bg"]
            bg_hex = re.sub(r"^linear-gradient\([^)]*\)$", colors["bg"], bg_hex)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = safe_hex(bg_hex, colors["bg"])

            # ── Title ──
            title_tag = slide_html.find(["h1", "h2"])
            if title_tag:
                t_box = slide.shapes.add_textbox(
                    safe_inch(0.8), safe_inch(0.5), 
                    safe_inch(11.7), safe_inch(1.0)
                )
                t_box.text_frame.word_wrap = True
                p = t_box.text_frame.paragraphs[0]
                p.text = title_tag.get_text().strip()
                p.font.size = safe_font_size(32)
                p.font.bold = True
                p.font.color.rgb = safe_hex(extract_styles(title_tag).get("color"), colors["title"])

            # ── Cards ──
            cards = slide_html.find_all("div", class_=re.compile(r"(card|glass-card)"))
            if cards:
                num_cards = min(len(cards), 4)  # Max 4 cards per slide
                card_w = max(1.0, (11.733 - (0.4 * (num_cards - 1))) / num_cards)
                
                for idx, card in enumerate(cards[:num_cards]):
                    c_styles = extract_styles(card)
                    left_pos = 0.8 + idx * (card_w + 0.4)
                    
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, 
                        safe_inch(left_pos), safe_inch(2.2),
                        safe_inch(card_w), safe_inch(4.2)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = safe_hex(
                        c_styles.get("background") or c_styles.get("background-color"), 
                        colors["card_bg"]
                    )
                    
                    # Card content
                    tb = slide.shapes.add_textbox(
                        safe_inch(left_pos + 0.2), safe_inch(2.4),
                        safe_inch(card_w - 0.4), safe_inch(3.8)
                    )
                    tb.text_frame.word_wrap = True
                    
                    c_title = card.find(["h3", "h4"])
                    if c_title:
                        p = tb.text_frame.paragraphs[0]
                        p.text = c_title.get_text().strip()
                        p.font.size = safe_font_size(16)
                        p.font.color.rgb = safe_hex(extract_styles(c_title).get("color"), colors["accent"])
            
            # ── Tables ──
            table_tag = slide_html.find("table")
            if table_tag:
                rows = table_tag.find_all("tr")
                if rows and len(rows) <= 12:  # Max 12 rows
                    c_count = min(max(len(r.find_all(["th", "td"])) for r in rows), 6)  # Max 6 cols
                    tbl_shape = slide.shapes.add_table(
                        len(rows), c_count, 
                        safe_inch(0.8), safe_inch(2.0),
                        safe_inch(11.733), safe_inch(4.4)
                    )
                    tbl = tbl_shape.table
                    for r_idx, row in enumerate(rows):
                        cells = row.find_all(["th", "td"])[:c_count]
                        for c_idx, cell in enumerate(cells):
                            tbl.cell(r_idx, c_idx).text = cell.get_text().strip()
                            p = tbl.cell(r_idx, c_idx).text_frame.paragraphs[0]
                            p.font.size = safe_font_size(12)
                            p.font.color.rgb = safe_hex(None, colors["title"])
                            
                            is_header = cell.name == "th" or r_idx == 0
                            tbl.cell(r_idx, c_idx).fill.solid()
                            tbl.cell(r_idx, c_idx).fill.fore_color.rgb = safe_hex(
                                None, 
                                colors["table_header_bg"] if is_header else colors["table_row_bg"]
                            )
        except Exception as e:
            print(f"Skipped corrupt slide {i}: {str(e)}")
            continue  # Skip bad slides but save presentation

    prs.save(output)


# ─── Endpoints ────────────────────────────────────────────────────────────────

@app.get("/health")
def health():
    return {"status": "ok", "version": "1.1.0"}

@app.post("/export-pptx")
def export_pptx(req: ExportRequest, request: Request):
    _check_rate(request.client.host)
    
    if not req.html or "<section" not in req.html:
        raise HTTPException(400, "HTML must contain <section> slides")
    
    cleaned_html = sanitize_html(req.html)
    if len(cleaned_html) > MAX_HTML_SIZE:
        raise HTTPException(413, f"HTML too large ({len(cleaned_html)} bytes)")

    buffer = io.BytesIO()
    try:
        compile_html_to_pptx(cleaned_html, req.theme, buffer)
    except Exception as e:
        error_id = f"ERR-{int(time.time())}"
        print(f"[{error_id}] Compile failed: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(500, detail={
            "error": "Compile failed",
            "error_id": error_id,
            "detail": str(e),
            "suggestion": "Simplify slide content and retry"
        })
    
    buffer.seek(0)
    filename = req.filename if req.filename.endswith('.pptx') else f"{req.filename}.pptx"
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )
